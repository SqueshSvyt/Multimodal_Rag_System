import os
import logging
from typing import Dict, Any, List


class ComprehensiveFileReader:
    """
    A comprehensive class to read various file types including text, documents, images, and more.
    """

    def __init__(self):
        """
        Initialize the FileReader with logging and supported file types.
        """
        logging.basicConfig(level=logging.INFO)
        self.logger = logging.getLogger(__name__)

        # Supported file types and their handlers
        self.supported_extensions = {
            # Text files
            '.txt': self._read_text_file,
            '.html': self._read_text_file,
            '.css': self._read_text_file,
            '.js': self._read_text_file,
            '.json': self._read_json_file,
            '.xml': self._read_text_file,
            '.csv': self._read_text_file,
            '.md': self._read_text_file,
            '.py': self._read_text_file,
            '.java': self._read_text_file,
            '.cpp': self._read_text_file,
            '.c': self._read_text_file,
            '.sql': self._read_text_file,
            '.yml': self._read_text_file,
            '.yaml': self._read_text_file,

            # Document files
            '.pdf': self._read_pdf_file,
            '.docx': self._read_docx_file,

            # Spreadsheet files
            '.xlsx': self._read_excel_file,
            '.xls': self._read_excel_file,

            # Presentation files
            '.pptx': self._read_pptx_file,

            # Image files (OCR extraction)
            '.jpg': self._read_image_file,
            '.jpeg': self._read_image_file,
            '.png': self._read_image_file,
            '.gif': self._read_image_file,
            '.bmp': self._read_image_file,
            '.tiff': self._read_image_file,
            '.webp': self._read_image_file,

            # Archive files
            '.zip': self._read_zip_file,
        }

    def read_file(self, file_path: str) -> Dict[str, Any]:
        """
        Read any supported file type and return its content.

        Args:
            file_path: Path to the file.

        Returns:
            A dictionary containing the file content and metadata.
        """
        try:
            # Check if file exists
            if not os.path.exists(file_path):
                self.logger.error(f"File not found: {file_path}")
                return self._create_error_response(file_path, f"File not found: {file_path}")

            # Get file extension
            extension = os.path.splitext(file_path)[1].lower()

            # Check if file type is supported
            if extension not in self.supported_extensions:
                self.logger.warning(f"Unsupported file extension: {extension}")
                return self._create_error_response(
                    file_path,
                    f"Unsupported file extension: {extension}. Supported types: {list(self.supported_extensions.keys())}"
                )

            # Get file info
            file_size = os.path.getsize(file_path)
            file_name = os.path.basename(file_path)

            # Read file using appropriate handler
            handler = self.supported_extensions[extension]
            content = handler(file_path)

            if content is None:
                return self._create_error_response(file_path, "Failed to extract content")

            self.logger.info(f"Successfully read file: {file_path}")
            return {
                "file_path": file_path,
                "file_name": file_name,
                "file_size": file_size,
                "file_type": extension,
                "content": content,
                "success": True,
                "error": None
            }

        except Exception as e:
            self.logger.error(f"Error reading file {file_path}: {str(e)}")
            return self._create_error_response(file_path, str(e))

    def _create_error_response(self, file_path: str, error_msg: str) -> Dict[str, Any]:
        """Create standardized error response."""
        return {
            "file_path": file_path,
            "file_name": os.path.basename(file_path) if file_path else None,
            "file_size": 0,
            "file_type": None,
            "content": None,
            "success": False,
            "error": error_msg
        }

    # Text file handlers
    def _read_text_file(self, file_path: str) -> str:
        """Read plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            raise Exception("Could not decode file with any supported encoding")

    def _read_json_file(self, file_path: str) -> str:
        """Read JSON files."""
        import json
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, indent=2)
        except Exception as e:
            # Fallback to reading as text
            return self._read_text_file(file_path)

    # Document file handlers
    def _read_pdf_file(self, file_path: str) -> str:
        """Read PDF files."""
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                text = ""
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                return text
        except ImportError:
            raise Exception(
                "PDF libraries not installed. Install with: pip install PyPDF2 or pip install pdfplumber")

    def _read_docx_file(self, file_path: str) -> str:
        """Read DOCX files."""
        try:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except ImportError:
            raise Exception("python-docx not installed. Install with: pip install python-docx")

    # Spreadsheet handlers
    def _read_excel_file(self, file_path: str) -> str:
        """Read Excel files."""
        try:
            import pandas as pd
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            content = ""
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                content += f"Sheet: {sheet_name}\n"
                content += df.to_string() + "\n\n"
            return content
        except ImportError:
            raise Exception("pandas and openpyxl not installed. Install with: pip install pandas openpyxl")

    # Presentation handlers
    def _read_pptx_file(self, file_path: str) -> str:
        """Read PowerPoint files."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text
        except ImportError:
            raise Exception("python-pptx not installed. Install with: pip install python-pptx")

    # Image handlers
    def _read_image_file(self, file_path: str) -> str:
        """Read images using OCR."""
        try:
            from PIL import Image
            import pytesseract

            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)

            if text.strip():
                return f"OCR extracted text from image:\n{text}"
            else:
                return f"Image file: {os.path.basename(file_path)} (no text detected)"
        except ImportError:
            return f"Image file: {os.path.basename(file_path)} (OCR libraries not installed)"
        except Exception as e:
            return f"Image file: {os.path.basename(file_path)} (OCR failed: {str(e)})"

    # Archive handlers
    def _read_zip_file(self, file_path: str) -> str:
        """Read ZIP files and list contents."""
        try:
            import zipfile
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                file_list = zip_ref.namelist()
                content = f"ZIP Archive contents ({len(file_list)} files):\n"
                content += "\n".join(file_list)

                # Try to extract text from text files within the archive
                text_files = [f for f in file_list if any(f.lower().endswith(ext) for ext in ['.txt', '.md', '.json'])]
                if text_files:
                    content += "\n\nExtracted text files:\n"
                    for text_file in text_files[:5]:  # Limit to first 5 text files
                        try:
                            with zip_ref.open(text_file) as f:
                                file_content = f.read().decode('utf-8')
                                content += f"\n--- {text_file} ---\n{file_content[:1000]}...\n"
                        except Exception:
                            continue

                return content
        except Exception as e:
            return f"ZIP archive: {os.path.basename(file_path)} (could not read: {str(e)})"

    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return list(self.supported_extensions.keys())

    def is_supported(self, file_path: str) -> bool:
        """Check if file type is supported."""
        extension = os.path.splitext(file_path)[1].lower()
        return extension in self.supported_extensions